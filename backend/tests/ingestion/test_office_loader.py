"""T11: docx headings->section_heading; pptx slide no->anchor; xlsx sheet->anchor, row-wise
blocks (AC-16/AC-17/AC-18)."""

import docx
import openpyxl
import pytest
from pptx import Presentation

from app.core.settings import Settings
from app.ingestion.loaders.office import load_office


def _settings() -> Settings:
    return Settings(
        _env_file=None,
        DATABASE_URL="postgresql+asyncpg://u:p@localhost:5432/db",
        ADMIN_EMAIL="admin@example.com",
        ADMIN_PASSWORD="secret",
        OPENAI_API_KEY="sk-test",
        PINECONE_API_KEY="pc-test",
        PINECONE_INDEX="campus-rag",
    )


def _make_docx(path):
    d = docx.Document()
    d.add_paragraph("Introduction", style="Heading 1")
    d.add_paragraph("Intro paragraph text about PU regulations.")
    d.add_paragraph("Section One", style="Heading 2")
    d.add_paragraph("Section one content about probation.")
    d.save(str(path))


def _make_pptx(path):
    prs = Presentation()
    layout = prs.slide_layouts[1]  # title + content
    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "Welcome Slide"
    slide1.placeholders[1].text = "Body text for slide one."

    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "Second Slide"
    slide2.placeholders[1].text = "Body text for slide two."
    prs.save(str(path))


def _make_xlsx(path):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "FeeSchedule"
    ws.append(["Program", "Fee"])
    ws.append(["BS CS", 50000])
    ws.append(["BS EE", 48000])
    wb.save(str(path))


@pytest.mark.asyncio
async def test_docx_headings_become_section_heading(tmp_path):
    path = tmp_path / "sample.docx"
    _make_docx(path)

    blocks = await load_office(path, "doc-docx", _settings(), "docx")

    intro = next(b for b in blocks if "Intro paragraph" in b.page_content)
    assert intro.metadata["section_heading"] == "Introduction"
    assert intro.metadata["anchor"] == "introduction"  # every block also gets an anchor
    section = next(b for b in blocks if "probation" in b.page_content)
    assert section.metadata["section_heading"] == "Section One"
    assert section.metadata["anchor"] == "section-one"
    assert all(b.metadata["doc_id"] == "doc-docx" for b in blocks)
    assert all(b.metadata["anchor"] is not None for b in blocks)


@pytest.mark.asyncio
async def test_pptx_slide_number_becomes_anchor(tmp_path):
    path = tmp_path / "sample.pptx"
    _make_pptx(path)

    blocks = await load_office(path, "doc-pptx", _settings(), "pptx")

    slide1_block = next(b for b in blocks if "slide one" in b.page_content)
    assert slide1_block.metadata["anchor"] == "slide-1"
    assert slide1_block.metadata["section_heading"] == "Welcome Slide"

    slide2_block = next(b for b in blocks if "slide two" in b.page_content)
    assert slide2_block.metadata["anchor"] == "slide-2"
    assert slide2_block.metadata["section_heading"] == "Second Slide"


@pytest.mark.asyncio
async def test_xlsx_row_wise_blocks_with_sheet_anchor(tmp_path):
    path = tmp_path / "sample.xlsx"
    _make_xlsx(path)

    blocks = await load_office(path, "doc-xlsx", _settings(), "xlsx")

    assert len(blocks) == 3  # header row + 2 data rows
    assert all(b.metadata["anchor"] == "FeeSchedule" for b in blocks)
    assert "BS CS" in blocks[1].page_content
    assert "50000" in blocks[1].page_content


@pytest.mark.asyncio
async def test_load_office_unsupported_file_type_raises(tmp_path):
    path = tmp_path / "sample.docx"
    _make_docx(path)

    with pytest.raises(ValueError, match="unsupported file_type"):
        await load_office(path, "doc-x", _settings(), "csv")
