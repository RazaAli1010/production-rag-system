"""T19: deterministic generator for the committed per-type ingestion fixtures.

Run with `python tests/fixtures/ingestion/generate_fixtures.py` (from `backend/`) to regenerate.
Fixtures are small, synthetic, and committed to the repo so `tests/ingestion` runs without any
network access or proprietary source documents. A genuine legacy `.doc`/`.ppt` binary fixture is
intentionally *not* included — that OLE2 format can't be produced from pure Python without a
real Word/PowerPoint or LibreOffice install; `test_legacy_loader.py` covers that path with a
synthetic OLE2-signature byte fixture built inline instead.
"""

from __future__ import annotations

from pathlib import Path

import docx
import fitz
import openpyxl
from pptx import Presentation

FIXTURES_DIR = Path(__file__).resolve().parent


def make_digital_pdf(path: Path) -> None:
    doc = fitz.open()
    p1 = doc.new_page()
    p1.insert_text((72, 72), "PU Academic Regulations 2021")
    p1.insert_text((72, 100), "Chapter 1: Admission Requirements")
    p1.insert_text(
        (72, 130),
        "Every candidate seeking admission to a degree programme at the University of the\n"
        "Punjab must satisfy the minimum eligibility criteria prescribed by the relevant\n"
        "faculty board, including the required aggregate percentage in prior examinations.",
    )
    p2 = doc.new_page()
    p2.insert_text((72, 72), "Chapter 2: Academic Probation")
    p2.insert_text(
        (72, 100),
        "A student whose CGPA falls below 2.0 in any regular semester shall be placed on\n"
        "academic probation. A student on probation who fails to raise the CGPA to the\n"
        "required threshold within two consecutive semesters may face suspension.",
    )
    doc.save(str(path))
    doc.close()


def make_scanned_pdf(path: Path) -> None:
    doc = fitz.open()
    for _ in range(2):
        page = doc.new_page()
        # small, single-channel, uniform-color pixmap: a "blank scanned page" image, kept tiny
        # on purpose (this fixture only needs to trip the has-image/no-text detection rule).
        pix = fitz.Pixmap(fitz.csGRAY, fitz.IRect(0, 0, 20, 25))
        pix.set_rect(pix.irect, (235,))
        page.insert_image(fitz.Rect(72, 72, 472, 572), pixmap=pix)
    doc.save(str(path), deflate=True)
    doc.close()


def make_html(path: Path) -> None:
    html = """<!DOCTYPE html>
<html>
<head><title>HEC Plagiarism Policy</title></head>
<body>
<nav><ul><li><a href="/">Home</a></li><li><a href="/policies">Policies</a></li></ul></nav>
<header>Higher Education Commission of Pakistan</header>
<h1 id="overview">Overview</h1>
<p>This policy establishes the standards for identifying and penalizing plagiarism in
research submitted by students and faculty across affiliated institutions.</p>
<h2 id="definitions">Definitions</h2>
<p>Plagiarism is defined as the unacknowledged use of another person's ideas, words, or
work, presented as one's own, whether intentional or unintentional.</p>
<h2 id="penalties">Penalties</h2>
<p>Confirmed cases of plagiarism may result in penalties ranging from a formal warning to
outright expulsion, depending on the severity and recurrence of the violation.</p>
<footer>Copyright HEC. All rights reserved.</footer>
</body>
</html>
"""
    path.write_text(html, encoding="utf-8")


def make_docx(path: Path) -> None:
    d = docx.Document()
    d.add_paragraph("PU Examination Rules", style="Title")
    d.add_paragraph("Attendance Requirement", style="Heading 1")
    d.add_paragraph(
        "A student must maintain at least seventy-five percent attendance in each course "
        "to be eligible to sit the final examination for that course."
    )
    d.add_paragraph("Grading Scale", style="Heading 1")
    d.add_paragraph(
        "Final grades are awarded on a relative grading scale approved by the Academic "
        "Council, ranging from A to F, with corresponding grade points from 4.0 to 0.0."
    )
    d.save(str(path))


def make_pptx(path: Path) -> None:
    prs = Presentation()
    layout = prs.slide_layouts[1]

    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "Orientation Week Schedule"
    slide1.placeholders[1].text = (
        "Day 1: Campus tour and registration.\nDay 2: Faculty introductions and course selection."
    )

    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "Student Support Services"
    slide2.placeholders[1].text = (
        "Counseling office located in the Admin Block.\nFinancial aid desk open Mon-Fri."
    )
    prs.save(str(path))


def make_xlsx(path: Path) -> None:
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "FeeSchedule2026"
    ws.append(["Programme", "Semester Fee (PKR)", "Admission Fee (PKR)"])
    ws.append(["BS Computer Science", 65000, 15000])
    ws.append(["BS Electrical Engineering", 62000, 15000])
    ws.append(["BA English", 35000, 10000])
    wb.save(str(path))


def main() -> None:
    make_digital_pdf(FIXTURES_DIR / "digital.pdf")
    make_scanned_pdf(FIXTURES_DIR / "scanned.pdf")
    make_html(FIXTURES_DIR / "sample.html")
    make_docx(FIXTURES_DIR / "sample.docx")
    make_pptx(FIXTURES_DIR / "sample.pptx")
    make_xlsx(FIXTURES_DIR / "sample.xlsx")
    print(f"Fixtures written to {FIXTURES_DIR}")


if __name__ == "__main__":
    main()
