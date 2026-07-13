"""Office document loaders — docx/pptx/xlsx (T11, AC-16/AC-17/AC-18).

Implemented directly on `python-docx`/`python-pptx`/`openpyxl` rather than the `Unstructured*`
wrapper loaders (design.md §3), for the same reason `loaders/pdf.py`/`loaders/html.py` bypass
their wrappers: fine per-block control over heading/slide/sheet metadata that the flattened
wrapper output doesn't expose cleanly. These are the same libraries `unstructured` itself uses
under the hood for these formats.

Note the parameter order is `(path, doc_id, settings, file_type)` — `file_type` last, not third
as design.md §4's signature lists it — so `routing.select_loader` can bind it via
`functools.partial(load_office, file_type=file_type)` and still expose the uniform
`(path, doc_id, settings) -> list[Document]` `LoaderFn` shape every other loader has.
"""

from __future__ import annotations

from pathlib import Path

import anyio
from docx import Document as DocxDocument
from langchain_core.documents import Document
from openpyxl import load_workbook
from pptx import Presentation

from app.core.settings import Settings
from app.ingestion.loaders.html import slugify

_HEADING_STYLE_PREFIX = ("Heading", "Title")


def _load_docx_sync(path: Path, doc_id: str) -> list[Document]:
    """AC-16: DOCX headings (Word styles `Heading *`/`Title`) become `section_heading` for the
    paragraphs that follow, until the next heading.

    Also assigns `anchor` (a slug of the nearest heading, or `para-{n}` for content before any
    heading) — design.md §3's routing table lists DOCX `anchor` as always `None`, but that
    leaves every DOCX block without a page *or* an anchor, contradicting AC-24/the feature DoD's
    ">= 95% of blocks carry page or anchor" and making DOCX content uncitable. Mirrors the same
    heading-slug approach `loaders/html.py` already uses.
    """
    docx = DocxDocument(str(path))
    blocks: list[Document] = []
    current_heading: str | None = None
    current_anchor: str | None = None
    seen_slugs: set[str] = set()
    pre_heading_index = 0

    for para in docx.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style_name = (para.style.name if para.style else "") or ""
        if style_name.startswith(_HEADING_STYLE_PREFIX):
            anchor = slugify(text)
            base, i = anchor, 2
            while anchor in seen_slugs:
                anchor = f"{base}-{i}"
                i += 1
            seen_slugs.add(anchor)
            current_heading, current_anchor = text, anchor
            continue

        anchor = current_anchor
        if anchor is None:
            pre_heading_index += 1
            anchor = f"para-{pre_heading_index}"

        blocks.append(
            Document(
                page_content=text,
                metadata={"doc_id": doc_id, "anchor": anchor, "section_heading": current_heading},
            )
        )
    return blocks


def _load_pptx_sync(path: Path, doc_id: str) -> list[Document]:
    """AC-17: slide number -> `anchor` (`slide-{n}`), slide title -> `section_heading`."""
    prs = Presentation(str(path))
    blocks: list[Document] = []

    for slide_no, slide in enumerate(prs.slides, start=1):
        anchor = f"slide-{slide_no}"
        title_shape = slide.shapes.title
        title = None
        if title_shape is not None and title_shape.has_text_frame:
            title = title_shape.text_frame.text.strip() or None

        for shape in slide.shapes:
            if shape == title_shape or not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            blocks.append(
                Document(
                    page_content=text,
                    metadata={"doc_id": doc_id, "anchor": anchor, "section_heading": title},
                )
            )
    return blocks


def _load_xlsx_sync(path: Path, doc_id: str) -> list[Document]:
    """AC-18: sheet name -> `anchor`; one block per non-empty row (row-wise, good-enough for
    merged-cell/multi-row-header sheets per requirements.md §5 "out of scope")."""
    wb = load_workbook(str(path), read_only=True, data_only=True)
    blocks: list[Document] = []

    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if not cells:
                continue
            blocks.append(
                Document(
                    page_content=" | ".join(cells),
                    metadata={"doc_id": doc_id, "anchor": sheet.title, "section_heading": None},
                )
            )
    wb.close()
    return blocks


_LOADERS = {"docx": _load_docx_sync, "pptx": _load_pptx_sync, "xlsx": _load_xlsx_sync}


async def load_office(
    path: Path, doc_id: str, settings: Settings, file_type: str
) -> list[Document]:
    """CPU/IO parsing (`python-docx`/`python-pptx`/`openpyxl` are sync libraries) — threaded per
    the project async rule."""
    loader = _LOADERS.get(file_type)
    if loader is None:
        raise ValueError(f"load_office: unsupported file_type={file_type!r}")
    return await anyio.to_thread.run_sync(loader, path, doc_id)
